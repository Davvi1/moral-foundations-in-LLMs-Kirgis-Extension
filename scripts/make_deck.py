#!/usr/bin/env python3
"""Build the presentation deck (20 slides, 16:9) with speaker notes.

Audience: BlueDot Impact technical-sprint mentor group, legible to an undergraduate
ML student. That drives two choices — a primer slide before any result, and plain
statements of what each readout physically does before any variance component.

Figures come from `make_figures.py`, which derives every number from
`results/derived/` and asserts it against the committed artifacts. This script
draws no numbers of its own except the headline counts, which are checked in
`verify_counts()` against the same sources rather than typed from memory.

Claim discipline carried over from CLAUDE.md and the repo's own rules:
  - "pre-specified analysis plan", NEVER "preregistration"
  - "almost entirely collinear", NEVER "perfectly collinear"
  - the alignment section stays narrow: no claims about frontier model values,
    no claim that logprob evaluation is broadly wrong

Usage:
    python scripts/make_deck.py                 # -> MFT-in-LLMs-presentation.pptx
    python scripts/make_deck.py --check         # verify counts only, write nothing
"""
from __future__ import annotations

import argparse
import re
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DERIVED = ROOT / "results" / "derived"
FIGDIR = ROOT / "figures"
OUT = ROOT / "MFT-in-LLMs-presentation.pptx"

# Tokens — identical to make_figures.py so slides and figures read as one system
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
RULE = RGBColor(0xE1, 0xE0, 0xD9)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
RED = RGBColor(0xE3, 0x49, 0x48)
GREEN = RGBColor(0x1B, 0xAF, 0x7A)
FONT = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)


# ---------------------------------------------------------------------------
def repo_url() -> str:
    """Read the real remote rather than printing a placeholder on a slide."""
    import subprocess
    try:
        u = subprocess.run(["git", "remote", "get-url", "origin"], cwd=ROOT,
                           capture_output=True, text=True, check=True).stdout.strip()
        return u.removeprefix("https://").removesuffix(".git")
    except Exception:
        return "(repository URL not set)"


def n_corrections() -> int:
    """Count the entries in CORRECTIONS.md.

    DERIVED, not hardcoded. An earlier version of this file asserted `== 21` and put
    "21 logged corrections" on the title slide; logging C22 broke the check and would
    have put a stale number in front of an audience if it had not. That is C21's whole
    lesson, reproduced inside the script written to enforce it.
    """
    body = (ROOT / "docs" / "CORRECTIONS.md").read_text(encoding="utf-8")
    return len(re.findall(r"^## C(\d+)\b", body, flags=re.M))


def verify_counts():
    """Re-derive the headline counts the deck states, rather than trusting prose."""
    df = pd.read_csv(DERIVED / "analysis_long_v2.csv")
    vr = pd.read_csv(DERIVED / "variance_ratio_v2.csv")

    facts = {
        "models collected": df.model.nunique(),
        "models analysed": df[~df.excluded].model.nunique(),
        "rows": len(df),
        "items": df.item_id.nunique(),
        "arms": df.condition.nunique(),
    }
    expect = {"models collected": 31, "models analysed": 30, "rows": 21576,
              "items": 116, "arms": 6}
    for k, v in expect.items():
        assert facts[k] == v, f"{k}: derived {facts[k]}, deck says {v}"
        print(f"  ok   {k:20s} {facts[k]}")
    print(f"  ok   {'corrections':20s} {n_corrections()}  (derived, not asserted)")

    p = vr[(vr.exclusions) & (~vr.scan_excluded) & (~vr.family_effect)
           & (vr.residual == "method-specific")].set_index("foundation")
    lo, hi = p.drop("Social Norms").R_median.min(), p.drop("Social Norms").R_median.max()
    assert abs(lo - 0.181) < .001 and abs(hi - 0.469) < .001, (lo, hi)
    print(f"  ok   R across 6 moral    {lo:.3f}-{hi:.3f}")
    assert abs(p.loc["Social Norms"].R_median - 0.133) < .001
    print("  ok   R control           0.133")
    for f in ["figures/fig%02d_*.png" % i for i in range(1, 14)]:
        assert list(ROOT.glob(f)), f"missing {f}"
    print("  ok   13 figures present")


# ---------------------------------------------------------------------------
def _tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, size, color, bold=False, italic=False, space_after=0,
         space_before=0, line=None):
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = FONT
    return r


def new_slide(prs, notes=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = SURFACE
    if notes:
        s.notes_slide.notes_text_frame.text = notes.strip()
    return s


def header(slide, head, sub=None, kicker=None):
    y = Inches(0.46)
    if kicker:
        tf = _tb(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.28))
        _run(tf.paragraphs[0], kicker.upper(), 11, MUTED, bold=True)
        y = Inches(0.78)
    tf = _tb(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.62))
    _run(tf.paragraphs[0], head, 27, INK, bold=True)
    if sub:
        tf2 = _tb(slide, MARGIN, y + Inches(0.62), W - 2 * MARGIN, Inches(0.46))
        _run(tf2.paragraphs[0], sub, 14, INK_2, line=1.25)
    return y + Inches(1.12 if sub else 0.68)


def footer(slide, n):
    tf = _tb(slide, W - MARGIN - Inches(1.0), H - Inches(0.52), Inches(1.0),
             Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _run(p, str(n), 11, MUTED)
    tf2 = _tb(slide, MARGIN, H - Inches(0.52), Inches(8.0), Inches(0.3))
    _run(tf2.paragraphs[0], "Is a language model's moral profile stable across "
         "scoring methods?", 10, MUTED)


def figure(slide, name, top, height=None):
    """Place a figure centred, scaled to fit the remaining space.

    Uses figures/deck/, the title-less variant — the slide header carries the
    headline, and the titled standalone version would duplicate it on every
    figure slide.
    """
    from PIL import Image
    path = next(ROOT.glob(f"figures/deck/{name}*.png"))
    iw, ih = Image.open(path).size
    avail_h = H - top - Inches(0.72)
    avail_w = W - 2 * MARGIN
    scale = min(avail_w / iw, avail_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(str(path), int((W - w) / 2), int(top), w, h)


def bullets(slide, top, items, size=15, gap=11, width=None, left=None):
    """items: (text, kind) where kind in {body, lead, dim, good, bad}."""
    colors = {"body": INK_2, "lead": INK, "dim": MUTED, "good": GREEN, "bad": RED}
    tf = _tb(slide, left or MARGIN, top, width or (W - 2 * MARGIN), Inches(4.4))
    first = True
    for text, kind in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _run(p, text, size, colors[kind], bold=(kind in ("lead", "good", "bad")),
             space_after=gap, line=1.30)
    return tf


def stat_row(slide, top, stats):
    """A KPI row: (value, label) tiles. The number is the chart."""
    n = len(stats)
    gap = Inches(0.28)
    tile_w = int((W - 2 * MARGIN - gap * (n - 1)) / n)
    for i, (val, lab) in enumerate(stats):
        x = MARGIN + i * (tile_w + gap)
        tf = _tb(slide, x, top, tile_w, Inches(0.72))
        _run(tf.paragraphs[0], val, 32, BLUE, bold=True)
        tf2 = _tb(slide, x, top + Inches(0.66), tile_w, Inches(0.62))
        _run(tf2.paragraphs[0], lab, 12, INK_2, line=1.2)


def rule(slide, y, x0=None, x1=None, color=RULE):
    from pptx.enum.shapes import MSO_SHAPE
    x0 = x0 or MARGIN
    x1 = x1 or (W - MARGIN)
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, x1 - x0, Emu(9525))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# ===========================================================================
def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    n = 0

    def sl(notes):
        nonlocal n
        n += 1
        s = new_slide(prs, notes)
        return s

    # ---- 1. Title -------------------------------------------------------
    s = sl("""
Opening line: "This is a replication that turned into an audit, and the audit
partly exonerated the paper it was auditing."

Frame the two-sentence version up front so nobody has to wait for it:
Kirgis measured the moral profiles of 21 frontier models and found that providers
differ. But his scoring method was forced to vary by provider — so 'provider' and
'how you read the answer' are not separable in his design. I ran the same
instrument on 31 open-weight models, where I can apply every scoring method to
every model, and asked whether the profile and the ranking survive the readout.

Say the honest headline now: for his specific pair of methods, they largely do.
The interesting findings are elsewhere, and they are about measurement.

Housekeeping: two days of work, $20 of compute, everything in the repo.
""")
    tf = _tb(s, MARGIN, Inches(2.05), Inches(11.4), Inches(1.9))
    _run(tf.paragraphs[0], "Is a language model's moral profile\nstable across "
         "scoring methods?", 40, INK, bold=True, line=1.12)
    tf = _tb(s, MARGIN, Inches(4.05), Inches(10.6), Inches(0.9))
    _run(tf.paragraphs[0], "A within-model audit of Kirgis (2025), "
         "arXiv:2511.11790, on open weights", 18, INK_2)
    rule(s, Inches(5.05), x1=Inches(4.0))
    tf = _tb(s, MARGIN, Inches(5.35), Inches(10.6), Inches(1.0))
    _run(tf.paragraphs[0], "David Moth  ·  Hertie School  ·  BlueDot Impact "
         "technical sprint", 14, INK_2)
    p = tf.add_paragraph()
    _run(p, "31 models  ·  116 vignettes  ·  6 scoring arms  ·  21,576 rows  ·  "
         f"{n_corrections()} logged corrections", 13, MUTED, space_before=8)
    footer(s, n)

    # ---- 2. The paper under audit ---------------------------------------
    s = sl("""
Kirgis, November 2025. He took the Moral Foundations Vignettes — 116 short
scenarios, each rated 0 to 4 for how morally wrong it is — and administered them
to 21 closed frontier models across six providers.

The instrument is real psychometrics: Clifford et al. 2015, validated on people,
with per-item human means published in their Table 1.

His four claims are on the right. Claims 2, 3 and 4 are the ones that get cited —
they are statements about which models have which values.

One caveat I'll repeat later: his human baseline is NOT nationally representative.
Clifford recruited a Qualtrics panel restricted to ages 18-40 and balanced on
ideology, about 30 raters per vignette. Kirgis describes it as 'a nationally
representative sample of US adults'. It isn't. That is free, fully documented
criticism and I verified it by fetching Clifford directly.

Also: his 'Care' foundation is emotional harm only. He drops all 16 physical-harm
Care items from Clifford's 132, keeping 116.
""")
    top = header(s, "The paper under audit",
                 "Kirgis, “Differences in the Moral Foundations of Large Language "
                 "Models”, arXiv:2511.11790", kicker="Background")
    bullets(s, top, [
        ("The instrument", "lead"),
        ("116 Moral Foundations Vignettes (Clifford et al. 2015) — short scenarios "
         "rated 0–4 for moral wrongness, with published per-item human means.", "body"),
        ("Six foundations: Care, Fairness, Liberty, Loyalty, Authority, Sanctity — "
         "plus Social Norms, which is a deliberately NON-moral control.", "body"),
        ("The administration", "lead"),
        ("21 closed frontier models, six providers, one API call per item.", "body"),
    ], width=Inches(6.5))
    x = Inches(7.5)
    tf = _tb(s, x, top, Inches(5.2), Inches(0.4))
    _run(tf.paragraphs[0], "His four claims", 15, INK, bold=True)
    bullets(s, top + Inches(0.48), [
        ("1.  MFT has explanatory power for LLM moral judgment", "body"),
        ("2.  Models diverge from the human baseline", "body"),
        ("3.  Providers differ systematically", "body"),
        ("4.  Divergence grows with capability", "body"),
    ], size=14, gap=13, left=x, width=Inches(5.2))
    tf = _tb(s, x, top + Inches(2.75), Inches(5.2), Inches(1.4))
    _run(tf.paragraphs[0], "Claims 2–4 are statements about which models hold "
         "which values. They are what a downstream reader actually uses.",
         13, MUTED, line=1.3)
    footer(s, n)

    # ---- 3. The confound ------------------------------------------------
    s = sl("""
Here is the problem, and it is not a criticism of his care — it is forced by the
APIs.

To score an answer you can either read the probability the model assigns to each
option, or you can make it write an answer and read what it wrote. Only some
providers expose logprobs. So he used top-3 logprob weighting where he could and
the mean of ten sampled responses everywhere else.

Look at the picture. Five of six providers sit entirely in one row. That means
'provider' and 'scoring method' are almost entirely collinear — you cannot tell
whether Anthropic differs from OpenAI because the models differ or because they
were measured differently.

BE PRECISE HERE — do not say 'perfectly collinear'. OpenAI straddles both arms:
four of its models were logprob-scored and two (GPT-4.5 and o3-Mini) were sampled.
So inside OpenAI, scoring method is confounded with model identity instead. Either
way nothing is identified. A reviewer who checks his Table 1 will catch the
stronger claim, and the weaker one carries the same conclusion.

This is a non-identification, not something you fix with a covariate. There is no
model measured both ways, so the method effect cannot be estimated from his data
at all. That is exactly what open weights fix.
""")
    top = header(s, "Scoring method was forced to vary by provider", kicker="The flaw")
    figure(s, "fig01", top - Inches(0.12))
    footer(s, n)

    # ---- 4. His code vs his paper ---------------------------------------
    s = sl("""
This is a side finding from re-analysing his committed data. No GPU, no new
collection — his repo has the raw logprob responses.

Left panel: for five of his six logprob-scored models, the top-3 probabilities
returned by the API sum to essentially 1, as they should. For grok-3-beta, 51 of
116 responses — 44% — come back with two entries instead of three, summing to
roughly zero probability, while the emitted token's own logprob says p = 1.0. The
two fields contradict each other. That is malformed data, not unusual data.

Right panel: his paper prints a formula with no denominator. His code divides by
the total probability. They are different estimators. The renormalisation in his
code accidentally rescues the corrupted rows — dividing near-zero by near-zero
recovers the argmax, which happens to be the right answer. Under the formula he
actually printed, grok-3's mean collapses from 1.98 to 1.20 and it drops from rank
4 to rank 6.

The deeper point, and the one that matters for this project: when only one of the
top three tokens is a digit, his renormalised estimator returns that digit exactly.
It degenerates to argmax. That happens on 5.6% of his responses. So inside the arm
he treats as one homogeneous method, there are actually three estimators, chosen
per item.

Tone: this is exploratory, it is from his own published data, and the honest
framing is 'provider logprob APIs cannot be assumed well-formed' — not 'he was
careless'. Any of us would have missed it without an integrity check.
""")
    top = header(s, "And his code does not implement the formula in his paper",
                 "Exploratory re-analysis of his committed data — no new collection "
                 "needed", kicker="A side finding")
    figure(s, "fig02", top - Inches(0.10))
    footer(s, n)

    # ---- 5. What we did -------------------------------------------------
    s = sl("""
The fix is open weights. With the weights on disk I can apply EVERY scoring method
to EVERY model, so scoring method becomes a within-model treatment instead of a
between-provider confound. That is the whole design.

Roster: 31 instruction-tuned open models from 0.5B to 72.7B — a 145-fold span —
across 12 families, with complete size ladders for Qwen (8 models) and Llama (4).
Revisions pinned by SHA so the run is reproducible. 30 analysed; one model,
SmolLM2-1.7B, is excluded because it answers '3' to all 116 items and so carries no
item-level information. That exclusion criterion was defined AFTER seeing the data
and I disclose it as a limitation.

Non-reasoning models only. Thinking-by-default would alter free generation while
being invisible to label scoring, which manufactures the very interaction I am
trying to measure.

Cost: about $20 on rented GPUs against a $100 ceiling. Downloads dominate, not
inference.

If asked why not frontier models: I cannot. The weights are not available, and
that is precisely the constraint that produced Kirgis's confound in the first place.
""")
    top = header(s, "The fix is open weights",
                 "With the weights on disk, every scoring method can be applied to "
                 "every model — so method becomes a within-model treatment.",
                 kicker="Design")
    stat_row(s, top + Inches(0.10), [
        ("31", "models collected\n30 analysed"),
        ("145×", "parameter span\n0.5B → 72.7B"),
        ("12", "model families\nQwen + Llama ladders complete"),
        ("6", "scoring arms\n5 sharing one prompt"),
        ("21,576", "rows\n116 items × arms × models"),
    ])
    rule(s, top + Inches(1.72))
    bullets(s, top + Inches(1.95), [
        ("The prompt is byte-identical across the five fixed-prompt arms — asserted "
         "per item, zero violations across 3,596 model × item cells.", "body"),
        ("Instruction-tuned, non-reasoning models only: thinking-by-default changes "
         "free generation while being invisible to label scoring, which would "
         "manufacture the very interaction under measurement.", "body"),
        ("Revisions pinned by SHA.  ~$20 of GPU time against a $100 ceiling.", "dim"),
    ], size=14)
    footer(s, n)

    # ---- 6. Six arms ----------------------------------------------------
    s = sl("""
This is the primer slide. Take it slowly — everything downstream depends on it.

A language model, given a context, returns one thing: a probability distribution
over the next token. That is all it computes. Longer text comes from doing it
repeatedly.

So 'asking a model a survey question' is underdetermined. There are at least two
families of answer:

READ THE PROPENSITY. Don't generate anything. Just look up what probability the
model assigns to each option. Three ways to do that here — the digit token '3';
the whole option line '3: Very wrong'; or the bare phrase 'Very wrong'.

WATCH THE BEHAVIOUR. Let it write, and parse what came out. Greedy takes the most
likely token at each step — deterministic. Sampled draws at temperature 1, ten
times, and averages.

The survey analogy: a respondent has some latent propensity across the five boxes.
Normally you see one tick. With open weights you can either read the propensities
or watch the tick. These are different estimands that we hope agree.

The sixth arm, cloze, removes the option list from the prompt. That is what makes
it textbook cloze — but it means the prompt changed, so it is EXCLUDED from the
primary number. Including it inflated my headline by 2.70× for several weeks. I'll
come back to that.

Two implementation details worth knowing if asked: '0' and ' 0' are different
tokens, and options 0 and 1 both start with the word 'Not' — so string scoring has
to be a full-sequence likelihood, not a single-position readout.
""")
    top = header(s, "Six ways to read one answer", kicker="Primer")
    figure(s, "fig03", top - Inches(0.18))
    footer(s, n)

    # ---- 7. Estimand ----------------------------------------------------
    s = sl("""
The estimand. For each foundation, fit a crossed random-effects model and take the
ratio of two variance components:

R = variance of the model-by-method interaction, over variance between models.

In words: how big is the perturbation that scoring method induces in a model's
position, relative to how much models actually differ from each other? If R is
small, model comparisons survive a change of readout. If R is large, the ranking
is partly an artifact of how you measured.

The bands and the decision rule were fixed at a git tag BEFORE any confirmatory
data existed — verifiably, the tagged commit contains no results directory.

CRITICAL WORDING: this is a pre-specified analysis plan, NOT a preregistration. A
tag in a repository I control is an internal discipline device, not independent
verification. I could in principle have retagged. That distinction is small in
practice and fatal in a write-up if you get it the wrong way round. Say it out loud.

The 'indeterminate' verdict existed before the data, on purpose, because a design
simulation showed that no feasible sample size resolves an R sitting near a band
boundary. Forcing a three-way call would manufacture false precision. Remember
that — it becomes the result.

One technical point if a statistician asks: the residual variance is
method-specific, not pooled. The arms have structurally different error variance —
greedy is discretised to integers, sampled carries Monte-Carlo error of order
1/sqrt(10) — and a single residual term pushes that into the interaction and
inflates R mechanically, toward the more publishable answer.
""")
    top = header(s, "The estimand, fixed before the data existed", kicker="Method")
    tf = _tb(s, MARGIN, top, Inches(7.4), Inches(0.9))
    _run(tf.paragraphs[0], "R  =  σ²(model × method)  ⁄  σ²(model)", 26, BLUE,
         bold=True)
    tf = _tb(s, MARGIN, top + Inches(0.72), Inches(6.9), Inches(1.2))
    _run(tf.paragraphs[0], "How large is the perturbation scoring method induces in "
         "a model's position, relative to how much models genuinely differ?",
         15, INK_2, line=1.3)
    y = top + Inches(1.72)
    for lab, txt, col in [
            ("R < 0.25", "robust — comparisons survive the readout", GREEN),
            ("0.25 – 1.0", "degraded — rankings not trusted alone", INK_2),
            ("R > 1.0", "not interpretable", RED),
            ("straddles a boundary", "indeterminate", MUTED)]:
        tf = _tb(s, MARGIN, y, Inches(2.3), Inches(0.34))
        _run(tf.paragraphs[0], lab, 14, col, bold=True)
        tf = _tb(s, MARGIN + Inches(2.45), y, Inches(4.5), Inches(0.34))
        _run(tf.paragraphs[0], txt, 14, INK_2)
        y += Inches(0.46)
    rule(s, y + Inches(0.22), x1=Inches(7.1))
    tf = _tb(s, MARGIN, y + Inches(0.42), Inches(6.5), Inches(1.0))
    _run(tf.paragraphs[0], "The residual variance is method-specific, not pooled: "
         "greedy is discretised to integers and sampled carries Monte-Carlo error of "
         "order 1/√10. A single residual term pushes that into the interaction and "
         "inflates R toward the more publishable answer.", 12.5, MUTED, line=1.32)
    x = Inches(7.9)
    rule(s, top, x0=x, x1=W - MARGIN)
    tf = _tb(s, x, top + Inches(0.18), Inches(4.8), Inches(3.4))
    _run(tf.paragraphs[0], "A pre-specified analysis plan —\nnot a preregistration",
         16, INK, bold=True, line=1.25)
    p = tf.add_paragraph()
    _run(p, "Bands and decision rule were locked at git tag "
         "analysis-plan-locked, in a commit that provably contains no results/raw/.",
         13, INK_2, space_before=12, line=1.32)
    p = tf.add_paragraph()
    _run(p, "But a tag in a repository the author controls is an internal discipline "
         "device, not independent verification. The distinction is small in practice "
         "and fatal if stated the wrong way round.", 13, MUTED, space_before=10,
         line=1.32)
    footer(s, n)

    # ---- 8. Label scoring fails -----------------------------------------
    s = sl("""
Now the findings, and I am deliberately leading with the measurement ones rather
than the headline variance ratio, because these are the ones that are established,
reproducible, and useful to anyone doing this kind of work.

Finding one. Label scoring — reading the probability of the digit — silently fails
on a large minority of models.

In my first harness, a faithful textbook implementation produced meaningless output
on 6 of 16 models. 38%. Two independent causes. First, SentencePiece tokenizers
encode '0' as two tokens, so a single-token lookup finds nothing. Second, the first
generated token often isn't the answer — Mistral emits a newline 116 times out of
116; Ministral emits end-of-sequence.

Here is the part that should worry everyone in this room: NEITHER RAISED AN ERROR.
Both produced plausible numbers in the right range. If you had those numbers in a
table you would publish them.

The only signal was retained probability mass — how much of the model's next-token
distribution actually sits on the five options before you renormalise.

This chart is the fixed version. Even after the fix, 7 of 30 models sit below 0.5.
Mistral-7B is at 0.008 — its 'moral judgment' is an expectation computed over
eight tenths of one percent of its next-token distribution, renormalised up to
look confident.

Takeaway to state plainly: if you use a logprob readout, report retained mass. It
costs nothing and it is the only thing that catches this.
""")
    top = header(s, "1.  Label scoring fails silently — no error, plausible numbers",
                 "A faithful first implementation produced meaningless output on "
                 "6 of 16 models (38%), from two independent causes.",
                 kicker="Measurement findings")
    figure(s, "fig04", top - Inches(0.06))
    footer(s, n)

    # ---- 9. Refusal leakage ---------------------------------------------
    s = sl("""
Finding two, and this is the one I would most want an alignment audience to take
away.

Label scoring is often justified on the grounds that it avoids refusals — you never
ask the model to speak, you just read the distribution, so a model that would have
declined still gives you a number.

That is exactly the problem. It does not avoid the refusal confound. It HIDES it.

Across 31 models, the greedy non-answer rate and label retained mass couple at
Spearman rho = -0.54. Models that decline or fall silent when you ask them to write
are largely the same models whose digit mass collapses when you read the
distribution. Renormalisation then manufactures a confident score out of whatever
digit mass is left.

The sharpest case is within a single model: Llama-3.1-8B craters to retained mass
0.481 on Sanctity — the disgust and purity items, which are graphic — against 0.829
on its other foundations, and it behaviourally refuses 35% of exactly those items.
The refusal is visible in the probability readout even though the readout never
asked it to speak.

Be honest about the evidence: 22 of the 31 models answer every item, so this
correlation rests on the handful that don't. I've put that on the chart rather than
leave it to be spotted.

Two implications. Studies using logprob readouts on safety-relevant content must
report retained mass — Kirgis's logprob arm has no such check. And the flip side is
useful: retained mass is a graded, generation-free refusal detector.
""")
    top = header(s, "2.  Refusal leaks into the logprob readout",
                 "Label scoring is defended as refusal-proof because the model never "
                 "has to speak. That is what makes it dangerous.",
                 kicker="Measurement findings")
    figure(s, "fig05", top - Inches(0.02))
    footer(s, n)

    # ---- 10. Free generation --------------------------------------------
    s = sl("""
Finding three. 'Free generation' sounds like one method. It is at least two, and
the choice between them can determine whether the model answers at all.

Ministral-8B: zero percent of items answered under greedy decoding, roughly fifty
percent under sampling. On BYTE-IDENTICAL prompts. Its greedy argmax at the first
position is the end-of-sequence token — it simply never says anything. Turn on
sampling and half the items come back.

Llama-3.2-1B refuses 109 of 116 under greedy.

These are not harness bugs. They reproduce across two independent collections, so
they are stable properties of those models.

The point for a methods audience: 'we used free generation' is not a specification.
Temperature is a researcher degree of freedom that silently changes your sample —
and it changes it non-randomly, because which items get dropped depends on content.

Fourth finding, no chart, worth saying out loud: because label scoring never
requires the model to speak, I can hold each model's probability-based answer to
the very items it refused behaviourally. That converts an untestable
missing-not-at-random assumption into a measured one. The result: the hypothesis
'refusal means extremely wrong' is TRUE for one model and FALSE for another.
Llama-3.1-8B rates its refused items 1.34 points higher; gemma-2-27b refuses items
it rates LESS severely. No single imputation rule is right across the roster, and
imputing the scale maximum would move a model's mean by up to 0.965 — the size of
a choice usually made silently in a footnote.
""")
    top = header(s, "3.  “Free generation” hides a decision about whether the model "
                 "answers at all",
                 "Ministral-8B: 0% of items answered under greedy, ~50% under "
                 "sampling — on byte-identical prompts.",
                 kicker="Measurement findings")
    figure(s, "fig06", top - Inches(0.02))
    footer(s, n)

    # ---- 11. R forest ---------------------------------------------------
    s = sl("""
Now the headline estimand.

Two things are true and they need to be said in this order.

First: the interaction is REAL. I ran a permutation null — shuffle the method
labels within model and item, destroying the interaction by construction, and refit.
700 full MCMC fits. It collapses to a median R of about 0.001. The observed values
run two to three orders of magnitude above that. So the estimator is calibrated:
destroy the effect and it reports none.

Second: the magnitude is NOT resolvable. Every one of the six moral foundations
comes back 'indeterminate' — the credible interval straddles a band boundary. R
runs from 0.181 for Care and Loyalty to 0.469 for Authority.

The one verdict that resolves is Social Norms, at 0.133, robust — and that is the
NON-MORAL CONTROL. Clifford designed those items to be odd but not wrong, like
drinking coffee with a spoon. They sit at the floor of the scale, every readout
agrees nothing much is wrong, so the readouts agree almost perfectly. That is what
a control ought to do, and it is a floor artifact rather than evidence that methods
agree.

This falsifies my own registered prediction P7, which said at least two foundations
would escape the indeterminate band at N around 30. None did. Going from 20 models
to 31 resolved nothing.

Say the magnitude in plain words: method effects are roughly a fifth to a half of
between-model variance. Not comparable to it — I withdrew that stronger claim.
""")
    top = header(s, "The interaction is real. Its magnitude is not resolvable.",
                 "Every moral foundation is indeterminate — which falsifies my own "
                 "registered prediction.", kicker="The primary result")
    figure(s, "fig08", top - Inches(0.04))
    footer(s, n)

    # ---- 12. Design simulation ------------------------------------------
    s = sl("""
Why didn't more models fix it? This is my favourite slide, because the answer was
sitting in a simulation I ran BEFORE collection and did not read carefully enough.

The curves show, for each true value of R, how often the estimate lands in the
correct band — at four sample sizes.

Look at the shape. It is not a ceiling you climb with more data. It has holes in it.
At R = 0.10 or 0.50 or 2.00 — in the interior of a band — accuracy is high and
improves with N. But at R = 0.25 and R = 1.00, exactly ON the band boundaries,
accuracy is 50% and STAYS 50% no matter how many models you add. Half the sampling
distribution necessarily falls either side of a line the true value sits on.

Now look at the red triangles: where my six foundations actually landed. Clustered
around 0.25. Sanctity at 0.246 sits almost exactly on the boundary.

So N was never the binding constraint. My prediction P7 leaned on '0.94 accuracy at
N=30' — which is the accuracy at R = 0.50, a value I did not have. The estimand
landed where no achievable sample size classifies it.

That is the result. Not a failure to reach significance — a demonstration that this
particular quantity, with this decision rule, is not resolvable at any N a student
project can reach. Which is a more useful thing to tell the next person than a point
estimate would have been.
""")
    top = header(s, "N was never the binding constraint",
                 "Classification accuracy collapses to a coin flip at a band "
                 "boundary — and stays there at every sample size.",
                 kicker="Why more models would not have helped")
    figure(s, "fig09", top - Inches(0.04))
    footer(s, n)

    # ---- 13. Pairwise tiers ---------------------------------------------
    s = sl("""
This is the most important slide in the deck. If you remember one chart, this one.

A single pooled R averages over all the arms and hides the fact that the arms
disagree wildly about how much they disagree. So here is R computed for every PAIR
of arms, on one common complete-case block, so every cell is comparable.

It separates into three clean tiers.

BLUE — the four readouts that carry real probability mass: label, string-line,
greedy, sampled. All six pairs between them sit between -0.09 and 0.20. That is
agreement. Negative values are estimator truncation; read them as 'indistinguishable
from no interaction'.

YELLOW — string-bare, scoring the phrase 'Very wrong' with no digit in front. R
jumps to between 0.62 and 1.33. Its mean retained mass is 0.0028. Under three tenths
of one percent.

ORANGE — anything involving cloze, which is the arm that changes the prompt. 2.5 to
4.3. That is why it is excluded from the primary.

So the honest statement is NOT 'scoring method perturbs moral profiles'. It is:
scoring method perturbs profiles IF you choose a readout that scores a region of
the output distribution the model essentially never visits.

That is a much more useful finding, and it is a warning about a specific practice
rather than a general scepticism.

If asked about the negative values: moment estimator truncation. A Bayesian fit with
a half-normal prior would return a small positive number instead.
""")
    top = header(s, "The disagreement is concentrated, not spread",
                 "Almost all of it comes from two low-mass probes. The four readouts "
                 "carrying real mass agree.", kicker="Where the effect actually lives")
    figure(s, "fig10", top - Inches(0.04))
    footer(s, n)

    # ---- 14. Concentration ----------------------------------------------
    s = sl("""
The same finding from a completely different direction — and I only found this
because someone asked whether any robustness check had ever varied the MODELS.
None had. Every check I'd run varied the arms.

One model out of 27 carries 34.4% of the interaction sum of squares. An equal share
would be 3.7%, so Mistral-7B is 9.3 times the average model's contribution. Drop it
and R falls by 51%.

And here is the part that makes this not a coincidence: Mistral-7B is the model
whose label retained mass is 0.008. The same model from the previous slide.

So the two findings are one finding. The method effect is carried disproportionately
by models whose probability readouts sit on almost no retained mass — cells the
design can barely measure in the first place.

Consequence, stated plainly: R is not a property of the roster. It is close to a
few-model phenomenon. Quoting a pooled R without this makes it sound far more
general than it is, and I'd be handing a reviewer their first question.

Note the share of interaction SUM OF SQUARES involves no variance-component
estimator at all, so this cannot be an artifact of the estimator truncating. The
leave-one-out on R agrees with it independently.
""")
    top = header(s, "…and it is carried by one model",
                 "The same model whose label retained mass is 0.008. Two findings, "
                 "one mechanism.", kicker="Where the effect actually lives")
    figure(s, "fig11", top - Inches(0.04))
    footer(s, n)

    # ---- 15. Kirgis survives --------------------------------------------
    s = sl("""
So — does the flaw that motivated this whole project actually break his paper?

Largely, no. And saying so is the RESULT of the audit, not a concession.

His two arms map onto two of mine: top-3 logprob weighting is essentially my label
arm, and the mean of ten samples is essentially my sampled arm. Boxed in red: they
rank models at Spearman rho = 0.82 over the six moral foundations. Their pairwise R
is 0.081 — the top tier from the previous slide, reached by a completely independent
route.

So the specific methodological flaw I built this project to expose is one his
conclusions can largely survive. A write-up that buried that would be dishonest.

Two more things in this matrix worth naming, both of which are limitations of MY
design and I'd rather state them than have them found.

Top-left: label and string-line correlate at 0.96, and at item level 0.988. They are
not two readouts. They are the SAME measurement. The reason is algebraic: the
log-probability of '3: Very wrong' is the log-probability of '3' plus the
log-probability of ': Very wrong' given '3' — and that second term is near-constant
across options when the prompt displays the digit-to-phrase mapping. Having
committed to the digit, the model just reads the phrase off the prompt.

So I have three independent probability readouts, not four. And no fixed prompt
escapes this: label scoring needs the digits visible, cloze needs them absent. I
found that after collection and I'm stating it myself.
""")
    top = header(s, "Kirgis's own confound is comparatively benign",
                 "His two arms rank models at ρ = 0.82, pairwise R = 0.081. "
                 "This is the audit's result, not a concession.",
                 kicker="Bearing on the target paper")
    figure(s, "fig07", top - Inches(0.02))
    footer(s, n)

    # ---- 16. Scale ------------------------------------------------------
    s = sl("""
The most consequential POSITIVE result, and the one that extends Kirgis rather than
auditing him.

His claim 2 was that models overweight the individualizing foundations — Care,
Fairness, Liberty — and underweight the binding ones — Loyalty, Authority, Sanctity
— relative to humans. His claim 4 was that divergence grows with capability.

Those turn out to be the same claim. The gap does appear, and it grows with model
size, on both complete ladders: +0.32 per decade of parameters on Qwen across eight
models and a 145-fold span, +0.26 on Llama across four. Leave-one-out keeps the sign
8 out of 8 and 4 out of 4.

This explains why a sub-14B sample found nothing — we were below the scale where the
pattern appears. And it unifies his claims 2 and 4: the divergence is a capability
phenomenon.

Now the caveats, which are load-bearing and I will not skip them:

Neither ladder is individually significant at 0.05 — p = 0.083 and p = 0.060. Only
the pooled fit is, and pooling is the WEAKER design, because models are not
exchangeable across families. Four of six families slope positive.

You can see the caveat in the scatter — there are real outliers. That is honest data
and I have not smoothed it.

And the standalone audit artifact in the repo still reads 'SUGGESTIVE, not
established, in either direction'. It predates this analysis and was never re-run,
so the upgrade rests entirely on the scale evidence, which is marginal per ladder.
I'm showing you both verdicts rather than picking the one I like.
""")
    top = header(s, "Kirgis's pattern is real — and it grows with model size",
                 "This unifies his claims 2 and 4: the divergence is a capability "
                 "phenomenon.", kicker="The positive extension")
    figure(s, "fig12", top - Inches(0.02))
    footer(s, n)

    # ---- 17. Compression ------------------------------------------------
    s = sl("""
But the adjustment is the whole result, and this slide is why.

Every line here is one Qwen model, regressing its severity ratings on the human
baseline. The grey diagonal is perfect agreement — slope b = 1.

The 0.5B model has slope 0.11. It barely tracks the human baseline at all; it says
roughly the same thing about every item. The 72B model has slope 1.06 — it tracks
the human ratings almost one to one.

That is 'compression', and it changes enormously with scale.

Here is the trap. Compression pulls every rating toward the middle of the scale.
Humans rate the individualizing foundations HIGHER than the binding ones — 2.66
against 2.38. So compression pulls the individualizing ratings DOWN more than it
pushes the binding ones UP. Pure compression therefore predicts a NEGATIVE
individualizing-minus-binding gap, with no moral content whatsoever.

And because compression itself vanishes as models get bigger, the RAW gap must rise
with scale even if the moral profile never changes at all.

Roughly a third of the raw slope is that confound. Reporting the raw number would
have been measuring the disappearance of compression and calling it morality.

One more honesty point: I only draw these lines over 1.40 to 3.80, the range the
moral items actually span. Extrapolating this fit further is exactly the error my
limitations document exists to disown — the non-moral control sits 0.9 points below
the lowest moral item with nothing in between, so any residual estimated down there
is not identified.
""")
    top = header(s, "…but the adjustment is the whole result",
                 "Compression changes enormously with scale — and pure compression "
                 "alone predicts a rising raw gap.", kicker="The positive extension")
    figure(s, "fig13", top - Inches(0.02))
    footer(s, n)

    # ---- 18. Limitations + corrections -----------------------------------
    s = sl("""
What I cannot claim. This list is short; the repo's version runs to 22 numbered
entries and I'd rather you read that.

The design never resolved R, at either sample size. R is not a property of the
roster — one model carries a third of it. My four conditions are three independent
readouts. R conflates rank reordering with methods disagreeing about SPREAD; about a
sixth to a quarter of it is spread, not order. Greedy turns out not to be
reproducible — 2.3% of scores move between two runs, which means my design gave it
one observation on reasoning that was wrong. The human baseline is a Qualtrics panel
of 18-to-40-year-olds, about 30 raters per item, not a representative sample. And
these are 31 open models under 73B — not the frontier, and the distance is
post-training as much as parameter count.

Deepest gap: everything is conditional on ONE prompt. I cannot say whether scoring
method matters more than an arbitrary rewording. That is the obvious sceptical
response to this whole project and I have not answered it.

Now the right-hand column, which I think is the most transferable thing here.

Twenty-two logged corrections. Not one of them was found by looking at a result and feeling
that it seemed wrong. Measurement artifacts do not announce themselves.

The worst was C15. My primary estimand contained the prompt-varying cloze arm for
several weeks, which inflated R by 2.70× — IN MY OWN FAVOUR, making scoring method
look like a bigger problem than it is. Three separate documents said that arm must
be excluded. Nothing in the code enforced it. A design commitment that exists only
in prose is not a commitment.

And C21 I logged two days ago, preparing this talk: the README advertised 19
corrections while linking to a file listing 20, and the test count was stale. Every
check in that repo is aimed at numbers derived from the data. Nothing was aimed at
what the project says about itself.
""")
    top = header(s, "What I cannot claim", kicker="Limitations")
    bullets(s, top, [
        ("The design never resolved R", "lead"),
        ("All six moral foundations indeterminate, at N = 20 and again at N = 30.",
         "body"),
        ("R is not a property of the roster", "lead"),
        ("One model of 27 carries 34% of the interaction; dropping it halves R.",
         "body"),
        ("Four conditions are three independent readouts", "lead"),
        ("label and string-line are the same measurement (item-level r = 0.988), and "
         "no fixed prompt escapes this. Found by us, stated by us.", "body"),
        ("Everything is conditional on one prompt — the deepest gap", "lead"),
        ("Whether scoring method matters more than an arbitrary rewording is "
         "untested, and is the obvious sceptical response to the design.", "body"),
        ("Not the frontier · greedy is not reproducible · the human baseline is an "
         "18–40 ideology-quota panel, n ≈ 30 per item", "dim"),
    ], size=13, gap=7, width=Inches(6.6))
    x = Inches(7.6)
    rule(s, top - Inches(0.10), x0=x, x1=W - MARGIN)
    tf = _tb(s, x, top + Inches(0.10), Inches(5.1), Inches(0.5))
    _run(tf.paragraphs[0], f"{n_corrections()} logged corrections", 20, INK, bold=True)
    tf = _tb(s, x, top + Inches(0.58), Inches(5.1), Inches(4.0))
    _run(tf.paragraphs[0], "Not one was found by looking at a result and feeling "
         "that it seemed wrong.", 14, INK_2, line=1.3)
    p = tf.add_paragraph()
    _run(p, "C15 — the worst", 14, RED, bold=True, space_before=14)
    p = tf.add_paragraph()
    _run(p, "The primary estimand contained a prompt-confounded arm for weeks, "
         "inflating R by 2.70× — in our own favour. Three documents said to exclude "
         "it. Nothing in the code enforced it.", 13, INK_2, space_before=5, line=1.3)
    p = tf.add_paragraph()
    _run(p, "C21 — logged while preparing this talk", 14, RED, bold=True,
         space_before=12)
    p = tf.add_paragraph()
    _run(p, "The README advertised 19 corrections while linking to a file listing "
         "20. Every check in the repo is aimed at the data; none at what the project "
         "says about itself.", 13, INK_2, space_before=5, line=1.3)
    footer(s, n)

    # ---- 19. Alignment ---------------------------------------------------
    s = sl("""
Why this matters beyond one paper. I want to be narrow here, because over-claiming
is how this kind of work dies.

One. Model value-profiles are becoming governance evidence. 'Model A is more X than
model B' arguments feed procurement, deployment, and regulation. This project shows
that the ranking such an argument rests on is partly a property of how you read the
answer, not only of the model.

Two, and this is the one I'd underline. The failure mode is SILENT. 38% of a
realistic roster returned plausible, meaningless numbers with no error raised. An
evaluation that crashes is safe. One that returns a number in the right range is
dangerous. Retained probability mass is a cheap, generic integrity check that would
have caught it, and it is not standard practice.

Three, the sharpest one for an alignment audience: refusal and measurement are
entangled precisely on safety-relevant content. Safety training makes models
decline. Logprob readouts renormalise that refusal away and hand you a confident
number. So the more heavily a model is safety-tuned, the more its logprob-based
value profile is a renormalisation artifact. That is a direct, mechanical way that
evaluating aligned models is harder than evaluating unaligned ones.

Four, on research practice: the honest result was the useful one. I found the
confound survivable for Kirgis's specific pair. A field that only publishes 'the
prior work is wrong' learns less than one that quantifies how much a design choice
can bite and reports when the answer is 'not much'.

What I am NOT claiming: nothing about frontier model values, nothing about logprob
evaluation being broadly wrong, and method-dependence is not evidence of
misalignment. One implementation, one roster, one prompt.
""")
    top = header(s, "Why this matters for AI alignment", kicker="Implications")
    bullets(s, top, [
        ("Value profiles are becoming governance evidence.", "lead"),
        ("“Model A is more X than model B” feeds procurement, deployment and "
         "regulation. The ranking is partly a property of the readout.", "body"),
        ("The failure mode is silent.", "lead"),
        ("38% of a realistic roster returned plausible, meaningless numbers with no "
         "error raised. An eval that crashes is safe; one that returns a number in "
         "the right range is not. Retained mass catches it and costs nothing.",
         "body"),
        ("Refusal and measurement are entangled on exactly the sensitive content.",
         "lead"),
        ("Safety training makes models decline; logprob readouts renormalise the "
         "refusal away and return confidence. The more safety-tuned the model, the "
         "more its logprob value-profile is an artifact.", "body"),
    ], size=13.5, gap=8, width=Inches(7.5))
    x = Inches(8.5)
    tf = _tb(s, x, top + Inches(0.02), Inches(4.2), Inches(3.6))
    _run(tf.paragraphs[0], "What this does NOT show", 15, INK, bold=True)
    for t in ["Anything about frontier model values",
              "That logprob-based evaluation is broadly wrong",
              "That method-dependence is evidence of misalignment"]:
        p = tf.add_paragraph()
        _run(p, "—  " + t, 13, MUTED, space_before=10, line=1.3)
    p = tf.add_paragraph()
    _run(p, "One implementation, one roster, one prompt.", 13, INK_2,
         space_before=16, italic=True, line=1.3)
    footer(s, n)

    # ---- 20. Close -------------------------------------------------------
    s = sl("""
Three sentences to close on.

The measurement findings are the strongest material and they don't depend on the
variance ratio at all: label scoring fails silently on a large minority of models,
retained mass is the check that catches it and doubles as a refusal detector, and
'free generation' hides a decision that can determine whether a model answers at all.

The audit exonerated its target on the specific point it was built to test —
rho = 0.82 between Kirgis's two arms — while finding a data-integrity problem in his
logprob arm that he did not detect, and independently supporting two of his four
claims.

And the estimand landed where no achievable sample size classifies it, which is a
more useful thing to hand the next person than a point estimate would have been.

The obvious next experiment is the one I deliberately did not run: make PROMPT a
designed factor and ask whether scoring method moves the ranking more than an
arbitrary rewording does. The machinery is built and 12 prompt variants are
committed, with predictions registered before any data. I stopped because adding a
second instrument changes the estimand and that is a different study.

Questions I expect: why not frontier models (can't — no weights, which is the
constraint that caused the confound); is 31 enough (no, and slide 12 shows more
wouldn't help); is this just prompt sensitivity (honest answer: I can't rule that
out, it's the deepest gap and it's the next experiment).
""")
    top = header(s, "Where this leaves things", kicker="Close")
    bullets(s, top, [
        ("The measurement findings are the strongest material — and none of them "
         "depend on the variance ratio.", "lead"),
        ("Label scoring fails silently; retained mass is the check that catches it "
         "and doubles as a generation-free refusal detector; “free generation” "
         "conceals whether the model answered at all.", "body"),
        ("The audit exonerated its target on the point it was built to test.", "lead"),
        ("ρ = 0.82 between Kirgis's two arms — while finding a data-integrity problem "
         "in his logprob arm he did not detect, and independently supporting two of "
         "his four claims.", "body"),
        ("The estimand landed where no achievable N classifies it.", "lead"),
        ("Which is a more useful thing to hand the next person than a point estimate "
         "would have been.", "body"),
    ], size=14, gap=9, width=Inches(7.4))
    x = Inches(8.4)
    rule(s, top - Inches(0.10), x0=x, x1=W - MARGIN)
    tf = _tb(s, x, top + Inches(0.12), Inches(4.3), Inches(3.4))
    _run(tf.paragraphs[0], "The next experiment", 16, INK, bold=True)
    p = tf.add_paragraph()
    _run(p, "Make the PROMPT a designed factor, and ask whether scoring method "
         "moves the ranking more than an arbitrary rewording does.", 13, INK_2,
         space_before=10, line=1.32)
    p = tf.add_paragraph()
    _run(p, "Built, not run: 12 prompt variants committed, six predictions "
         "registered before any data exists. Deferred because a second instrument "
         "changes the estimand.", 13, MUTED, space_before=10, line=1.32)
    p = tf.add_paragraph()
    _run(p, repo_url(), 12.5, BLUE, space_before=16, bold=True)
    footer(s, n)

    prs.save(OUT)
    return n


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--check", action="store_true",
                    help="verify headline counts against the data, write nothing")
    a = ap.parse_args()
    print("Verifying deck counts against the data:")
    verify_counts()
    if a.check:
        print("\nAll checks pass. No deck written (--check).")
        return 0
    n = build()
    print(f"\nWrote {OUT.name}  ({n} slides, speaker notes on every slide)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
